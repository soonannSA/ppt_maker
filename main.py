from fastapi import FastAPI, Response
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64
from typing import List, Dict, Optional, Any
import os
import sys
import json
import zipfile
from PIL import Image
import random
import string

app = FastAPI()

# Ephemeral storage for cloud sync
sync_sessions: Dict[str, Any] = {}

def generate_session_id(length=6):
    return ''.join(random.choices(string.ascii_uppercase + string.digits, k=length))

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class ImageCrop(BaseModel):
    x: float
    y: float
    w: float
    h: float

class AppImage(BaseModel):
    src: str
    crop: ImageCrop

class CellStyles(BaseModel):
    backgroundColor: str
    color: str
    fontWeight: str
    fontSize: str
    fontStyle: str
    fontFamily: str
    textAlign: str
    imagePadding: float = 0.2

class Cell(BaseModel):
    text: str
    rowSpan: int = 1
    colSpan: int = 1
    hidden: bool = False
    image: Optional[AppImage] = None
    styles: CellStyles

class Row(BaseModel):
    height: float

class Col(BaseModel):
    width: float

class TableData(BaseModel):
    id: str
    x: float
    y: float
    width: float
    height: float
    gridColor: str = '#c9c9c9'
    rows: List[Row]
    cols: List[Col]
    cells: Dict[str, Cell]

class SlideData(BaseModel):
    id: str
    title: str = "APP"
    tables: List[TableData]

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    if hex_color == 'transparent' or not hex_color:
        return RGBColor(255, 255, 255)
    try:
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except:
        return RGBColor(255, 255, 255)



@app.post("/export")
def export_pptx(slides_data: List[SlideData]):
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    
    title_only_layout = prs.slide_layouts[5]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(title_only_layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
        
        for data in slide_data.tables:
            rows = len(data.rows)
            cols = len(data.cols)
            
            left = Cm(data.x)
            top = Cm(data.y)
            width = Cm(data.width)
            height = Cm(data.height)
            
            shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = shape.table

            # Set row heights
            for i, row in enumerate(data.rows):
                table.rows[i].height = Cm(row.height)
                
            # Set col widths
            for i, col in enumerate(data.cols):
                table.columns[i].width = Cm(col.width)
                
            # Track cell positions for image overlay
            col_offsets = [0.0]
            curr_x = 0.0
            for col in data.cols:
                curr_x += col.width
                col_offsets.append(curr_x)
                
            row_offsets = [0.0]
            curr_y = 0.0
            for row in data.rows:
                curr_y += row.height
                row_offsets.append(curr_y)

            # Apply Cell Processing
            for r in range(rows):
                for c in range(cols):
                    key = f"{r},{c}"
                    if key not in data.cells:
                        continue
                    
                    cell_data = data.cells[key]
                    p_cell = table.cell(r, c)
                    
                    # Handle Merges
                    if (cell_data.rowSpan > 1 or cell_data.colSpan > 1) and not cell_data.hidden:
                        try:
                            target_row = min(r + cell_data.rowSpan - 1, rows - 1)
                            target_col = min(c + cell_data.colSpan - 1, cols - 1)
                            other_cell = table.cell(target_row, target_col)
                            p_cell.merge(other_cell)
                        except Exception as e:
                            print(f"Merge failed: {e}")
                    
                    if cell_data.hidden:
                        continue
                    
                    # Text & Styling
                    p_cell.text = cell_data.text
                    if p_cell.text_frame.paragraphs:
                        paragraph = p_cell.text_frame.paragraphs[0]
                        if cell_data.styles.textAlign == 'center':
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif cell_data.styles.textAlign == 'right':
                            paragraph.alignment = PP_ALIGN.RIGHT
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT
                            
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.bold = cell_data.styles.fontWeight == 'bold'
                        run.font.italic = cell_data.styles.fontStyle == 'italic'
                        try:
                            fs_pt = float(cell_data.styles.fontSize.replace('px', '')) * 0.75
                            run.font.size = Pt(fs_pt)
                        except:
                            run.font.size = Pt(14)
                        run.font.name = cell_data.styles.fontFamily
                        run.font.color.rgb = hex_to_rgb(cell_data.styles.color)
                    
                    # Background Color
                    try:
                        if cell_data.styles.backgroundColor != 'transparent':
                            fill = p_cell.fill
                            fill.solid()
                            fill.fore_color.rgb = hex_to_rgb(cell_data.styles.backgroundColor)
                    except:
                        pass

                    # Image Overlay with Cropping / Contain
                    if cell_data.image and cell_data.image.src:
                        try:
                            img_data = cell_data.image.src
                            if ',' in img_data:
                                img_data = img_data.split(',')[1]
                            
                            image_bytes = base64.b64decode(img_data)
                            image_stream = io.BytesIO(image_bytes)
                            
                            # Aspect ratio logic
                            pil_img = Image.open(image_stream)
                            orig_w, orig_h = pil_img.size
                            image_stream.seek(0)
                            
                            cell_w_cm = sum(col.width for col in data.cols[c : c + cell_data.colSpan])
                            cell_h_cm = sum(row.height for row in data.rows[r : r + cell_data.rowSpan])
                            
                            padding_cm = cell_data.styles.imagePadding if hasattr(cell_data.styles, 'imagePadding') else 0.2
                            
                            avail_w_cm = max(0.1, cell_w_cm - 2 * padding_cm)
                            avail_h_cm = max(0.1, cell_h_cm - 2 * padding_cm)
                            
                            # Calculate aspect ratio
                            ratio_w = avail_w_cm / orig_w
                            ratio_h = avail_h_cm / orig_h
                            scale = min(ratio_w, ratio_h)
                            
                            final_w_cm = orig_w * scale
                            final_h_cm = orig_h * scale
                            
                            img_left = Cm(data.x + col_offsets[c] + padding_cm + (avail_w_cm - final_w_cm) / 2)
                            img_top = Cm(data.y + row_offsets[r] + padding_cm + (avail_h_cm - final_h_cm) / 2)
                            img_width = Cm(final_w_cm)
                            img_height = Cm(final_h_cm)
                            
                            picture = slide.shapes.add_picture(image_stream, img_left, img_top, img_width, img_height)
                            
                            # For 'contain', we ignore frontend crop unless it's strictly a custom crop the user explicitly did.
                            # The frontend now uses `object-fit: contain`, so we shouldn't apply full-cell crop.
                            # But if the user *manually* cropped it, should we apply it?
                            # For now, let's keep aspect ratio and NOT crop, as requested: "does conform to the aspect ration"
                            
                        except Exception as e:
                            print(f"Image overlay failed: {e}")

    # Save PPTX to bytes first
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)

    # Embed JSON state as a properly-registered custom part inside the ZIP.
    try:
        state_dict = {"slides": [s.model_dump() if hasattr(s, "model_dump") else s.dict() for s in slides_data]}
        state_json = json.dumps(state_dict).encode("utf-8")

        APP_STATE_PATH = "app_state.json"
        CONTENT_TYPE_ENTRY = (
            '<Override PartName="/app_state.json" '
            'ContentType="application/json"/>'
        )

        original_zip_bytes = pptx_io.read()
        output_io = io.BytesIO()

        with zipfile.ZipFile(io.BytesIO(original_zip_bytes), 'r') as zin:
            with zipfile.ZipFile(output_io, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        ct_xml = data.decode("utf-8")
                        if APP_STATE_PATH not in ct_xml:
                            ct_xml = ct_xml.replace(
                                "</Types>",
                                f"{CONTENT_TYPE_ENTRY}</Types>"
                            )
                        data = ct_xml.encode("utf-8")
                    zout.writestr(item, data)
                zout.writestr(APP_STATE_PATH, state_json)

        output_io.seek(0)
        final_bytes = output_io.read()
    except Exception as e:
        print(f"Error embedding state into ZIP: {e}")
        pptx_io.seek(0)
        final_bytes = pptx_io.read()

    # Robust Response: Set Content-Length and use standard DISPOSITION
    return Response(
        content=final_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": 'attachment; filename="table_bundle.pptx"',
            "Content-Length": str(len(final_bytes)),
            "Cache-Control": "no-cache"
        }
    )

@app.post("/sync/save")
def sync_save(state: Any):
    sid = generate_session_id()
    sync_sessions[sid] = state
    # Limit storage to 100 sessions to prevent memory leak
    if len(sync_sessions) > 100:
        first_key = next(iter(sync_sessions))
        del sync_sessions[first_key]
    return {"status": "ok", "session_id": sid}

@app.get("/sync/load/{session_id}")
def sync_load(session_id: str):
    if session_id in sync_sessions:
        return {"status": "ok", "data": sync_sessions[session_id]}
    return {"status": "error", "message": "Session not found or expired"}

@app.post("/import")
async def import_pptx(file: UploadFile = File(...)):
    try:
        content = await file.read()
        # Try new ZIP-based method first
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if 'app_state.json' in zf.namelist():
                    state_json = zf.read('app_state.json').decode('utf-8')
                    return {"status": "ok", "data": json.loads(state_json)}
        except Exception:
            pass
        # Fallback: check old XML comment method for backwards compatibility
        try:
            prs = Presentation(io.BytesIO(content))
            comments = prs.core_properties.comments
            if comments and comments.startswith("ppt_table_maker:"):
                state_json = comments[len("ppt_table_maker:"):]
                return {"status": "ok", "data": json.loads(state_json)}
        except Exception:
            pass
        return {"status": "error", "message": "No valid state found in this PPTX file."}
    except Exception as e:
        return {"status": "error", "message": f"Failed to read PPTX: {str(e)}"}

if getattr(sys, 'frozen', False):
    frontend_dir = os.path.join(sys._MEIPASS, "dist")
else:
    frontend_dir = os.path.join(os.path.dirname(__file__), "frontend_dist")

if os.path.exists(frontend_dir):
    app.mount("/", StaticFiles(directory=frontend_dir, html=True), name="static")
else:
    @app.get("/")
    def read_root():
        # Diagnostic info for Render deployment issues
        try:
            dir_contents = os.listdir(os.getcwd())
        except:
            dir_contents = ["Could not list directory"]
            
        return {
            "error": "Frontend build not found",
            "looking_at": frontend_dir,
            "directory_contents": dir_contents,
            "help": "Ensure the 'frontend_dist' folder is in the same folder as main.py on GitHub."
        }

if __name__ == "__main__":
    import uvicorn
    import multiprocessing
    multiprocessing.freeze_support()
    uvicorn.run(app, host="127.0.0.1", port=8000)
